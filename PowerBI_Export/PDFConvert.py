from pptx import Presentation
from pptx.util import Inches, Pt
import fitz  # PyMuPDF
import io
import os
import tkinter as tk
from tkinter import filedialog, ttk

def select_files(title="Select files", filetypes=[('All files', '*.*')]):
    root = tk.Tk()
    root.withdraw()  # Hide the main window
    file_paths = filedialog.askopenfilenames(title=title, filetypes=filetypes)
    root.destroy()
    return root.tk.splitlist(file_paths)

def select_file(title="Select a file", filetypes=[('All files', '*.*')]):
    root = tk.Tk()
    root.withdraw()  # Hide the main window
    file_path = filedialog.askopenfilename(title=title, filetypes=filetypes)
    root.destroy()
    return file_path

def select_folder(title="Select a folder"):
    root = tk.Tk()
    root.withdraw()  # Hide the main window
    folder_path = filedialog.askdirectory(title=title)
    root.destroy()
    return folder_path

def update_progress(progress_bar, value):
    progress_bar['value'] = value
    progress_bar.update()

# Select the PDF files and PowerPoint template
paths_to_pdfs = select_files("Select the PDF files", [('PDF files', '*.pdf')])
path_to_template = select_file("Select the PowerPoint template", [('PowerPoint files', '*.pptx')])

# Select the output folder
output_folder = select_folder("Select folder to save the PowerPoint")

# Check if the files were selected
if not paths_to_pdfs or not path_to_template or not output_folder:
    raise Exception("File not selected")

# Load the PowerPoint template
prs = Presentation(path_to_template)

# Define the slide dimensions (16:9 PowerPoint slide dimensions)
slide_height = Inches(7.5)  # inches
slide_width = Inches(13.33)  # inches

# Define the margins for the image
left_margin = (Inches(13.33) - slide_width) / 2  # inches
top_margin = (Inches(7.5) - slide_height) / 2  # inches

# Initialize progress bar window
progress_window = tk.Tk()
progress_window.title("Conversion Progress")
total_pages = sum([len(fitz.open(pdf_path)) for pdf_path in paths_to_pdfs])
progress_bar = ttk.Progressbar(progress_window, orient='horizontal', length=300, mode='determinate', maximum=total_pages)
progress_bar.pack()
progress_window.update()

# Process each PDF
for path_to_pdf in paths_to_pdfs:
    pdf_base_name = os.path.splitext(os.path.basename(path_to_pdf))[0]
    output_file = os.path.join(output_folder, pdf_base_name + ".pptx")

    # Load the PDF
    pdf_document = fitz.open(path_to_pdf)

    # Loop through the pages
    for page_num in range(len(pdf_document)):
        update_progress(progress_bar, page_num)

        # Convert the page to an image
        page = pdf_document.load_page(page_num)
        pix = page.get_pixmap(dpi=700)  # High resolution
        img_data = pix.tobytes("png")  # Convert the image to PNG bytes
        image_stream = io.BytesIO(img_data)  # Create a stream to add to the slide

        # Add a slide for the image
        slide_layout = prs.slide_layouts[6]  # Assuming '6' is the index for a blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Add the image to the slide
        pic = slide.shapes.add_picture(image_stream, left_margin, top_margin, slide_width, slide_height)

    # Save the presentation for the current PDF
    prs.save(output_file)
    print(f"Presentation for {pdf_base_name} saved at {output_file}")

# Final update to the progress bar
update_progress(progress_bar, total_pages)
progress_window.destroy()


