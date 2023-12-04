from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tkinter as tk
from tkinter import filedialog, ttk
import os
from threading import Thread

# Constants
EMU_PER_INCH = 914400

def crop_image(shape, left, right, top, bottom):
    shape.crop_left = left
    shape.crop_right = right
    shape.crop_top = top
    shape.crop_bottom = bottom

def expand_image_to_overflow(shape):
    # Calculate the increase in dimensions in EMUs
    width_increase_emu = int(0.03 * EMU_PER_INCH)
    height_increase_emu = int(0.03 * EMU_PER_INCH)
    # Adjust the width and height of the shape
    shape.width += width_increase_emu
    shape.height += height_increase_emu
    # Calculate position adjustment in EMUs
    left_shift_emu = int(0.01 * EMU_PER_INCH)
    top_shift_emu = int(0.01 * EMU_PER_INCH)
    # Adjust the position
    shape.left -= left_shift_emu
    shape.top -= top_shift_emu

def select_files(title="Select files"):
    root = tk.Tk()
    root.withdraw()
    file_paths = filedialog.askopenfilenames(title=title, filetypes=[('PowerPoint files', '*.pptx')])
    root.destroy()
    return file_paths

def select_save_folder(title="Select a folder"):
    root = tk.Tk()
    root.withdraw()
    folder_path = filedialog.askdirectory(title=title)
    root.destroy()
    return folder_path

# Function to process PowerPoint files (to be run in a separate thread)
def process_ppt_files(ppt_paths, save_path, update_ui_callback):
    for index, ppt_path in enumerate(ppt_paths):
        file_name = os.path.basename(ppt_path)
        prs = Presentation(ppt_path)
        crop_left_value = 0.020
        crop_right_value = 0.020
        crop_top_value = 0.035
        crop_bottom_value = 0.035

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    crop_image(shape, crop_left_value, crop_right_value, crop_top_value, crop_bottom_value)
                    expand_image_to_overflow(shape)

        save_path_with_filename = os.path.join(save_path, file_name)
        prs.save(save_path_with_filename)

        # Update the UI after processing each file
        update_ui_callback(index + 1)

# Function to update the UI
def update_ui(progress_bar, status_label, current, total_files):
    progress_bar['value'] = (current / total_files) * 100
    status_label.config(text=f'Processing {current}/{total_files} files...')
    if current == total_files:
        status_label.config(text="Done! Close this window to exit.")
        ttk.Button(progress_window, text="Close", command=progress_window.destroy).pack()

# Main function
def main():
    ppt_paths = select_files("Select PowerPoint files to adjust")
    total_files = len(ppt_paths)
    save_path = select_save_folder("Select folder to save adjusted PowerPoints")

    # Set up the progress bar window immediately
    global progress_window
    progress_window = tk.Tk()
    progress_window.title("Processing PowerPoints")
    ttk.Label(progress_window, text="Progress:").pack()
    progress_bar = ttk.Progressbar(progress_window, orient='horizontal', length=300, mode='determinate')
    progress_bar.pack()
    status_label = ttk.Label(progress_window, text="Starting...")
    status_label.pack()

    # Start processing the PPT files in a separate thread to keep UI responsive
    processing_thread = Thread(target=process_ppt_files, args=(ppt_paths, save_path, lambda current: update_ui(progress_bar, status_label, current, total_files)))
    processing_thread.start()

    progress_window.mainloop()

if __name__ == "__main__":
    main()


